"""Stage 3 — Document Parsing.

Per-type parser router: dispatches to the correct parser based on the file
category from Stage 1. Each parser extracts text and structural content
as faithfully as the format allows.

Key design decisions from the architecture doc:
- Native PDF → PyMuPDF/pdfplumber direct extraction (solved problem locally)
- Mixed PDF → page-by-page routing (native pages → text layer, scanned → OCR chain)
- Office → python-docx/pptx/openpyxl (no API call needed for structured formats)
- Text formats → direct parse (already machine-readable)
"""

from __future__ import annotations

import csv
import io
import json
import logging
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF

from src.models.schemas import (
    FileCategory,
    FileDetectionResult,
    PageContent,
    PageStructure,
    ParsedDocument,
    ClassificationResult,
)

logger = logging.getLogger(__name__)


def parse_document(
    file_path: str | Path,
    detection: FileDetectionResult,
    classification: ClassificationResult,
) -> ParsedDocument:
    """Route to the correct parser based on file category.

    Returns a ParsedDocument with per-page content, ready for OCR (Stage 4)
    on pages that need it.
    """
    path = Path(file_path)
    category = detection.file_category

    parsers = {
        FileCategory.PDF: _parse_pdf,
        FileCategory.DOCX: _parse_docx,
        FileCategory.DOC: _parse_doc,
        FileCategory.PPTX: _parse_pptx,
        FileCategory.PPT: _parse_ppt,
        FileCategory.XLSX: _parse_xlsx,
        FileCategory.XLS: _parse_xlsx,  # openpyxl handles both via xlrd fallback
        FileCategory.CSV: _parse_csv,
        FileCategory.TSV: _parse_tsv,
        FileCategory.MARKDOWN: _parse_text,
        FileCategory.PLAINTEXT: _parse_text,
        FileCategory.HTML: _parse_html,
        FileCategory.XML: _parse_text,
        FileCategory.JSON: _parse_json,
        FileCategory.IMAGE: _parse_image,
    }

    parser_fn = parsers.get(category)
    if parser_fn is None:
        logger.warning("No parser for category %s — returning empty document", category.value)
        return ParsedDocument(file_path=str(path), file_category=category)

    doc = parser_fn(path, classification)
    doc.file_path = str(path)
    doc.file_category = category

    logger.info(
        "Parsed '%s': %d pages, %d chars total",
        path.name,
        doc.total_pages,
        sum(len(p.text) for p in doc.pages),
    )
    return doc


# ---------------------------------------------------------------------------
# PDF parsing — the most complex case
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse a PDF, using per-page structural classification to decide approach."""
    doc = fitz.open(str(path))
    pages: list[PageContent] = []
    warnings: list[str] = []

    page_structures = classification.page_structures

    for page_num, page in enumerate(doc):
        # Get per-page structure (or default to overall)
        if page_num < len(page_structures):
            structure = page_structures[page_num]
        else:
            structure = classification.structural

        if structure == PageStructure.NATIVE_TEXT:
            # Extract text directly — this is the 95% case for native PDFs
            text = page.get_text("text")

            # Also extract as dict for structured blocks (used in Stage 5)
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)

            page_content = PageContent(
                page_number=page_num + 1,
                structure=PageStructure.NATIVE_TEXT,
                text=text,
                ocr_confidence=1.0,  # Native text → perfect confidence
                ocr_method="pymupdf_native",
            )

        elif structure == PageStructure.SCANNED:
            # Text will be filled by Stage 4 (OCR)
            page_content = PageContent(
                page_number=page_num + 1,
                structure=PageStructure.SCANNED,
                text="",  # Will be populated by OCR
                ocr_confidence=0.0,
                ocr_method="pending_ocr",
            )

        else:  # MIXED
            # Try native text extraction — if quality is low, mark for OCR
            text = page.get_text("text")
            if len(text.strip()) > 100:
                page_content = PageContent(
                    page_number=page_num + 1,
                    structure=PageStructure.MIXED,
                    text=text,
                    ocr_confidence=0.8,
                    ocr_method="pymupdf_native_mixed",
                )
            else:
                page_content = PageContent(
                    page_number=page_num + 1,
                    structure=PageStructure.SCANNED,
                    text="",
                    ocr_confidence=0.0,
                    ocr_method="pending_ocr",
                )

        pages.append(page_content)

    # Extract document-level metadata
    metadata = _extract_pdf_metadata(doc)
    doc.close()

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.PDF,
        total_pages=len(pages),
        pages=pages,
        metadata=metadata,
        warnings=warnings,
    )


def _extract_pdf_metadata(doc: fitz.Document) -> dict[str, Any]:
    """Extract PDF metadata (title, author, creation date, etc.)."""
    meta = doc.metadata or {}
    return {
        "title": meta.get("title", ""),
        "author": meta.get("author", ""),
        "subject": meta.get("subject", ""),
        "creator": meta.get("creator", ""),
        "producer": meta.get("producer", ""),
        "creation_date": meta.get("creationDate", ""),
        "modification_date": meta.get("modDate", ""),
        "page_count": doc.page_count,
    }


# ---------------------------------------------------------------------------
# Office format parsers
# ---------------------------------------------------------------------------

def _parse_docx(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse DOCX using python-docx."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)

    # Extract tables
    table_texts: list[str] = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        table_texts.append("\n".join(rows))

    if table_texts:
        full_text += "\n\n" + "\n\n".join(table_texts)

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.DOCX,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text=full_text,
                ocr_confidence=1.0,
                ocr_method="python_docx",
            )
        ],
    )


def _parse_doc(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse legacy DOC — try LibreOffice conversion to DOCX first."""
    import subprocess
    import tempfile

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, str(path)],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0:
                docx_path = Path(tmpdir) / f"{path.stem}.docx"
                if docx_path.exists():
                    return _parse_docx(docx_path, classification)

        logger.warning("LibreOffice conversion failed for '%s'", path.name)
    except FileNotFoundError:
        logger.warning("LibreOffice not installed — cannot parse .doc files")
    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice conversion timed out for '%s'", path.name)

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.DOC,
        total_pages=0,
        warnings=["Failed to parse .doc file — LibreOffice conversion required"],
    )


def _parse_pptx(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse PPTX — extract text from slides + speaker notes."""
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[PageContent] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Speaker Notes]: {notes}")

        page_content = PageContent(
            page_number=slide_num,
            structure=PageStructure.NATIVE_TEXT,
            text="\n".join(texts),
            ocr_confidence=1.0,
            ocr_method="python_pptx",
        )
        pages.append(page_content)

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.PPTX,
        total_pages=len(pages),
        pages=pages,
    )


def _parse_ppt(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse legacy PPT — convert via LibreOffice."""
    import subprocess
    import tempfile

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, str(path)],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0:
                pptx_path = Path(tmpdir) / f"{path.stem}.pptx"
                if pptx_path.exists():
                    return _parse_pptx(pptx_path, classification)
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        logger.warning("LibreOffice conversion failed for '%s': %s", path.name, e)

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.PPT,
        total_pages=0,
        warnings=["Failed to parse .ppt file — LibreOffice conversion required"],
    )


def _parse_xlsx(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse XLSX/XLS — each sheet becomes a page."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[PageContent] = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        rows: list[str] = []

        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))

        page_content = PageContent(
            page_number=sheet_num,
            structure=PageStructure.NATIVE_TEXT,
            text=f"Sheet: {sheet_name}\n" + "\n".join(rows),
            ocr_confidence=1.0,
            ocr_method="openpyxl",
        )
        pages.append(page_content)

    wb.close()

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.XLSX,
        total_pages=len(pages),
        pages=pages,
    )


def _parse_csv(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse CSV."""
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.CSV,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text="\n".join(rows),
                ocr_confidence=1.0,
                ocr_method="csv_reader",
            )
        ],
    )


def _parse_tsv(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse TSV."""
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter="\t")
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.TSV,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text="\n".join(rows),
                ocr_confidence=1.0,
                ocr_method="tsv_reader",
            )
        ],
    )


# ---------------------------------------------------------------------------
# Text-based format parsers
# ---------------------------------------------------------------------------

def _parse_text(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse plain text, Markdown, or XML — already machine-readable."""
    text = path.read_text(encoding="utf-8", errors="replace")

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.PLAINTEXT,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text=text,
                ocr_confidence=1.0,
                ocr_method="direct_read",
            )
        ],
    )


def _parse_html(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse HTML — extract text content, strip tags."""
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.texts: list[str] = []
            self._skip_tags = {"script", "style", "head"}
            self._in_skip = 0

        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            if tag.lower() in self._skip_tags:
                self._in_skip += 1

        def handle_endtag(self, tag: str) -> None:
            if tag.lower() in self._skip_tags:
                self._in_skip = max(0, self._in_skip - 1)

        def handle_data(self, data: str) -> None:
            if self._in_skip == 0 and data.strip():
                self.texts.append(data.strip())

    html_content = path.read_text(encoding="utf-8", errors="replace")
    extractor = _TextExtractor()
    extractor.feed(html_content)
    text = "\n".join(extractor.texts)

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.HTML,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text=text,
                ocr_confidence=1.0,
                ocr_method="html_parser",
            )
        ],
    )


def _parse_json(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse JSON — pretty-print for downstream processing."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        text = raw  # Malformed JSON — pass through as-is

    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.JSON,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.NATIVE_TEXT,
                text=text,
                ocr_confidence=1.0,
                ocr_method="json_parser",
            )
        ],
    )


def _parse_image(path: Path, classification: ClassificationResult) -> ParsedDocument:
    """Parse standalone image — mark for OCR processing in Stage 4."""
    return ParsedDocument(
        file_path=str(path),
        file_category=FileCategory.IMAGE,
        total_pages=1,
        pages=[
            PageContent(
                page_number=1,
                structure=PageStructure.SCANNED,
                text="",  # Will be filled by OCR
                ocr_confidence=0.0,
                ocr_method="pending_ocr",
            )
        ],
    )
